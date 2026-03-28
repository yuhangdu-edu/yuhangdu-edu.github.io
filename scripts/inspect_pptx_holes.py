"""
Read the PPTX file and print all shape geometries for slide 2 (Layer 1 nuli).
Find the 6 Rounded Rectangle holes and the disc chart.
"""
from pptx import Presentation
from pptx.util import Emu
import math

PPTX = r'C:\Users\ydu\OneDrive - London Business School\0 PhD Life\Job Market\personal website\others\tool design\tool_design_updated.pptx'

prs = Presentation(PPTX)

# Slide dimensions
sw = prs.slide_width    # in EMU
sh = prs.slide_height   # in EMU
print(f"Slide size: {sw} x {sh} EMU  = {sw/914400:.2f}\" x {sh/914400:.2f}\"")
print(f"            {sw/360000:.0f} x {sh/360000:.0f} hundredths-of-inch")

slide = prs.slides[1]   # Slide 2 = Layer 1 nuli (0-indexed)
print(f"\nSlide 2 shapes ({len(slide.shapes)}):")

for shape in slide.shapes:
    l = shape.left   # EMU from left
    t = shape.top    # EMU from top
    w = shape.width  # EMU
    h = shape.height # EMU

    # Fractional positions relative to slide
    l_frac = l / sw
    t_frac = t / sh
    w_frac = w / sw
    h_frac = h / sh

    # Center position
    cx_frac = (l + w/2) / sw
    cy_frac = (t + h/2) / sh

    # Size in inches
    w_in = w / 914400
    h_in = h / 914400

    print(f"  [{shape.shape_id:3d}] {shape.name[:30]:30s}  "
          f"left={l_frac:.4f} top={t_frac:.4f} w={w_frac:.4f} h={h_frac:.4f}  "
          f"cx={cx_frac:.4f} cy={cy_frac:.4f}  "
          f"w_in={w_in:.3f}\" h_in={h_in:.3f}\"")

# Also print slide 4 (Layer 1 multi) for comparison
print("\nSlide 4 shapes (Layer 1 multi):")
slide4 = prs.slides[3]
for shape in slide4.shapes:
    l = shape.left; t = shape.top; w = shape.width; h = shape.height
    cx_frac = (l + w/2) / sw; cy_frac = (t + h/2) / sh
    w_in = w / 914400; h_in = h / 914400
    print(f"  [{shape.shape_id:3d}] {shape.name[:30]:30s}  "
          f"cx={cx_frac:.4f} cy={cy_frac:.4f} "
          f"w={w/sw:.4f} h={h/sh:.4f}  "
          f"w_in={w_in:.3f}\" h_in={h_in:.3f}\"")
